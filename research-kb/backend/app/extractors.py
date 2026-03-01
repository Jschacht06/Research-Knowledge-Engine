from pathlib import Path
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation


def extract_text(filepath: str) -> str:
    path = Path(filepath)
    ext = path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(filepath)
    if ext == ".docx":
        return _extract_docx(filepath)
    if ext == ".pptx":
        return _extract_pptx(filepath)

    return ""


def _extract_pdf(filepath: str) -> str:
    doc = fitz.open(filepath)
    parts = []
    for page in doc:
        parts.append(page.get_text("text"))
    doc.close()
    return "\n".join(parts).strip()


def _extract_docx(filepath: str) -> str:
    d = DocxDocument(filepath)
    parts = [p.text for p in d.paragraphs if p.text]
    return "\n".join(parts).strip()


def _extract_pptx(filepath: str) -> str:
    prs = Presentation(filepath)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts).strip()